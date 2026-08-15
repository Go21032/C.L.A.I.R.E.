"""PDF/Word/Excel/PowerPoint、およびPython/JSON/CSV等のテキスト/コード系ファイルを
テキスト抽出し、ノート取り込み(ingest_notes.py)と同じ経路(chunker→embed→LanceDB)で
「ナレッジ」として永続登録する。

11日目ノート(サポートAI作製計画/11日目Web検索対応・UIデザイン確定・マルチモーダル対応調査.md)④の
改善策のうち、実装コストの低いPDF/Word/Excel/PowerPoint対応(画像は対象外。同ノートに追記した
「画像対応の検討事項」参照)。voice_gateway.pyの`POST/GET /documents`・`DELETE /documents/{filename}`
から呼ばれる想定(CLIからも単体で使える)。

12日目追記: 📷ボタンを📎に統合したのに合わせ、TEXT_EXTENSIONS(.py/.json/.csv/.md等)を
追加し、コード/テキスト系ファイルもそのままナレッジ登録できるようにした(画像は引き続き
このモジュールの対象外。静的index.htmlの📎ボタン側でMIMEタイプ判定して振り分ける)。

source規約: `f"doc:{filename}"`(ingest_notes.pyの`f"note:{path}"`・memory_store.append_turnの
`f"chat:{chat_id}"`に倣う)。`role="document"` / `route="DOCUMENT"`で登録する。同名ファイルの
再アップロードは ingest_notes.py と同様に「既存行を削除してから追加」でupsert扱いにする。

前提(未インストールなら):
    pip install pdfplumber python-docx openpyxl python-pptx

使い方(CLI):
    python doc_ingest.py <file>              # 単発の取り込み
    python doc_ingest.py --list              # 登録済みナレッジの一覧
    python doc_ingest.py --delete <filename> # 指定ファイル名のナレッジを削除
"""
from __future__ import annotations

import argparse
import io
import uuid
from datetime import datetime
from pathlib import Path

import chunker

# memory_store は config.yaml(db_path等)をimport時に読み込むため、ingest_notes.py と同じ理由で
# ここではトップレベルimportせず、実際にDBへ触る関数の内部で遅延importする。

OFFICE_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx"}

# 12日目追記: 📷ボタン統合(画像は別経路)に合わせて、📎から
# コード/テキスト系ファイルも添付できるようにした。抽出はバイト列を
# デコードするだけ(見出し等は付けず、そのまま1ファイル分のテキストとして扱う)。
TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".csv", ".tsv", ".json", ".yaml", ".yml",
    ".py", ".js", ".jsx", ".ts", ".tsx", ".java", ".c", ".h", ".cpp", ".hpp",
    ".cs", ".go", ".rs", ".rb", ".php", ".sh", ".ps1", ".sql", ".html", ".css",
    ".xml", ".ini", ".toml", ".log",
}

SUPPORTED_EXTENSIONS = OFFICE_EXTENSIONS | TEXT_EXTENSIONS
DOC_SOURCE_PREFIX = "doc:"


class UnsupportedFileTypeError(ValueError):
    """対応していない拡張子のファイルが渡された。"""


def extract_text(filename: str, data: bytes) -> str:
    """拡張子に応じてファイルの中身をテキストへ変換する(見出し風の区切り付き)。"""
    suffix = Path(filename).suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf(data)
    if suffix == ".docx":
        return _extract_docx(data)
    if suffix == ".xlsx":
        return _extract_xlsx(data)
    if suffix == ".pptx":
        return _extract_pptx(data)
    if suffix in TEXT_EXTENSIONS:
        return _extract_plain_text(data)
    raise UnsupportedFileTypeError(
        f"未対応のファイル形式です: {suffix or '(拡張子なし)'}(対応: {sorted(SUPPORTED_EXTENSIONS)})"
    )


def _extract_plain_text(data: bytes) -> str:
    """テキスト/コード系ファイルをデコードする。UTF-8を優先し、Windowsのメモ帳等で
    保存されたShift_JIS(CP932)系のファイルにもフォールバックする。どちらでも
    デコードできない場合は文字化けを許容してでも中身を落とさない(errors="replace")。"""
    for encoding in ("utf-8-sig", "utf-8", "cp932"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _extract_pdf(data: bytes) -> str:
    import pdfplumber  # 遅延import(重い依存のためモジュールimport時のコストを避ける)

    pages = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = (page.extract_text() or "").strip()
            if text:
                pages.append(f"## page {i}\n{text}")
    return "\n\n".join(pages)


def _extract_docx(data: bytes) -> str:
    import docx  # python-docx

    document = docx.Document(io.BytesIO(data))
    parts = [p.text.strip() for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_xlsx(data: bytes) -> str:
    import openpyxl

    workbook = openpyxl.load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    parts: list[str] = []
    for sheet in workbook.worksheets:
        sheet_lines = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                sheet_lines.append(" | ".join(cells))
        if sheet_lines:
            parts.append(f"## sheet: {sheet.title}\n" + "\n".join(sheet_lines))
    return "\n\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    parts = []
    for i, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                texts.append(shape.text_frame.text.strip())
        if texts:
            parts.append(f"## slide {i}\n" + "\n".join(texts))
    return "\n\n".join(parts)


def _chunks_for_text(text: str) -> list[str]:
    return [c.strip() for c in chunker.chunk_markdown(text) if c.strip()]


def ingest_document(filename: str, data: bytes, *, text: str | None = None) -> dict:
    """アップロードされたファイルを抽出→チャンク化→LanceDBへ永続登録する(upsert)。

    text: 呼び出し元(voice_gateway.pyのupload_document())が、レスポンスへ含める
    かどうか(DOC_INLINE_MAX_CHARS判定)のために先にextract_text()を呼んでいる場合、
    ここで二重に抽出しないよう渡せる。Noneのまま(=従来どおりの呼び出し方)なら
    ここでextract_text()を呼ぶ(後方互換)。

    戻り値: {"filename": str, "chunks": int}
    """
    import memory_store  # 遅延import(理由は冒頭コメント参照)

    if text is None:
        text = extract_text(filename, data)
    chunks = _chunks_for_text(text)

    source = f"{DOC_SOURCE_PREFIX}{filename}"
    safe_source = source.replace("'", "''")
    table = memory_store._table()
    # 冪等性: 同一ファイル名の既存行を先に削除してから追加する(ingest_notes.pyと同じupsert方式)
    table.delete(f"source = '{safe_source}'")

    if chunks:
        now = datetime.now().isoformat(timespec="seconds")
        rows = [
            {
                "id": str(uuid.uuid4()),
                "date": now,
                "source": source,
                "role": "document",
                "route": "DOCUMENT",
                "topic": filename,
                "content": c,
                "vector": memory_store.embed(c, is_query=False),
            }
            for c in chunks
        ]
        table.add(rows)

    return {"filename": filename, "chunks": len(chunks)}


def list_documents() -> list[dict]:
    """登録済みナレッジをファイル名ごとに集計して返す(チャンク数・最終登録日時)。"""
    import memory_store  # 遅延import

    table = memory_store._table()
    if table.count_rows() == 0:
        return []
    df = table.to_pandas()
    docs = df[df["source"].str.startswith(DOC_SOURCE_PREFIX)]
    if docs.empty:
        return []
    grouped = docs.groupby("source").agg(chunks=("id", "count"), date=("date", "max")).reset_index()
    results = [
        {
            "filename": row["source"][len(DOC_SOURCE_PREFIX):],
            "chunks": int(row["chunks"]),
            "date": row["date"],
        }
        for _, row in grouped.iterrows()
    ]
    return sorted(results, key=lambda r: r["date"], reverse=True)


def get_document_text(filename: str) -> str | None:
    """登録済みファイルの抽出全文を、チャンクを連結して復元する。

    14日目①: ピン留めしたファイルの全文をクライアントが取り直すためのI/F。
    ingest_document()はチャンク化してLanceDBへ登録するだけで、全文を読み戻す
    経路がなかった(ノート⓪-3参照)。ここではノート内の検討どおり、まず
    「チャンクを連結して復元する」(方針1)で実装する。チャンクはingest_document()内で
    1回のtable.add()にまとめて追加されるため、to_pandas()の行順は概ね挿入順を保つ。
    オーバーラップ(chunk_markdown内)による重複が実害になった場合は、方針2
    (アップロード時に抽出全文を別途保存する)へ切り替える。

    未登録のファイル名にはNoneを返す(呼び出し元はこれを404として扱う)。
    """
    import memory_store  # 遅延import

    table = memory_store._table()
    if table.count_rows() == 0:
        return None
    df = table.to_pandas()
    source = f"{DOC_SOURCE_PREFIX}{filename}"
    rows = df[df["source"] == source]
    if rows.empty:
        return None
    return "\n\n".join(rows["content"].tolist())


def delete_document(filename: str) -> int:
    """指定ファイル名のナレッジを削除する。削除した行数を返す。"""
    import memory_store  # 遅延import

    source = f"{DOC_SOURCE_PREFIX}{filename}"
    safe_source = source.replace("'", "''")
    table = memory_store._table()
    before = table.count_rows()
    table.delete(f"source = '{safe_source}'")
    after = table.count_rows()
    return before - after


def main() -> None:
    parser = argparse.ArgumentParser(description="PDF/Word/Excel/PowerPointをナレッジ(記憶DB)へ取り込む")
    parser.add_argument("file", nargs="?", default=None, help="取り込むファイルのパス")
    parser.add_argument("--list", action="store_true", help="登録済みナレッジの一覧を表示する")
    parser.add_argument("--delete", metavar="FILENAME", default=None, help="指定ファイル名のナレッジを削除する")
    args = parser.parse_args()

    if args.list:
        for doc in list_documents():
            print(f"{doc['filename']}: {doc['chunks']}チャンク (最終登録: {doc['date']})")
        return

    if args.delete:
        deleted = delete_document(args.delete)
        print(f"{args.delete}: {deleted}行削除しました")
        return

    if not args.file:
        raise SystemExit("ファイルパス、--list、--delete のいずれかを指定してください")

    path = Path(args.file)
    if not path.exists():
        raise SystemExit(f"対象ファイルが存在しません: {path}")

    result = ingest_document(path.name, path.read_bytes())
    print(f"{result['filename']}: {result['chunks']}チャンクを登録しました")


if __name__ == "__main__":
    main()
