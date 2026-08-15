"""tests/test_doc_ingest.py
------------------------------
11日目ノート④: PDF/Word/Excel/PowerPointのテキスト抽出(doc_ingest.py)のテスト。

memory_store(config.yaml依存)へは触れない範囲(拡張子ディスパッチ・実際の
テキスト抽出そのもの)だけを対象にする。ingest_document()/list_documents()/
delete_document()はLanceDB接続が要るため、ここでは対象外(実機での動作確認は
smoke_test_pipe.py同様に手動で行う)。

PDFは書き込み用ライブラリ(reportlab等)をこの取り込み機能のためだけに追加する
コストを避け、実際のpdfplumber読み取りは手動確認に留める。拡張子ディスパッチ
(_extract_pdfが呼ばれること)自体はモックで検証する。
"""
from __future__ import annotations

import io
import sys
import unittest
from pathlib import Path
from unittest.mock import patch

TESTS_DIR = Path(__file__).resolve().parent
SCRIPTS_DIR = TESTS_DIR.parent
RAG_MEMORY_DIR = SCRIPTS_DIR / "rag_memory"
for p in (RAG_MEMORY_DIR, TESTS_DIR):
    if str(p) not in sys.path:
        sys.path.insert(0, str(p))

import doc_ingest  # noqa: E402


class TestExtractTextDispatch(unittest.TestCase):
    """拡張子ごとに正しい抽出関数へディスパッチされることを確認する。"""

    def test_unsupported_extension_raises(self):
        with self.assertRaises(doc_ingest.UnsupportedFileTypeError):
            doc_ingest.extract_text("image.png", b"dummy")

    def test_no_extension_raises(self):
        with self.assertRaises(doc_ingest.UnsupportedFileTypeError):
            doc_ingest.extract_text("noext", b"dummy")

    def test_pdf_dispatches_to_pdf_extractor(self):
        with patch.object(doc_ingest, "_extract_pdf", return_value="pdf-text") as mocked:
            result = doc_ingest.extract_text("a.PDF", b"data")  # 大文字拡張子でも動くこと
        mocked.assert_called_once_with(b"data")
        self.assertEqual(result, "pdf-text")

    def test_docx_dispatches_to_docx_extractor(self):
        with patch.object(doc_ingest, "_extract_docx", return_value="docx-text") as mocked:
            result = doc_ingest.extract_text("a.docx", b"data")
        mocked.assert_called_once_with(b"data")
        self.assertEqual(result, "docx-text")

    def test_xlsx_dispatches_to_xlsx_extractor(self):
        with patch.object(doc_ingest, "_extract_xlsx", return_value="xlsx-text") as mocked:
            result = doc_ingest.extract_text("a.xlsx", b"data")
        mocked.assert_called_once_with(b"data")
        self.assertEqual(result, "xlsx-text")

    def test_pptx_dispatches_to_pptx_extractor(self):
        with patch.object(doc_ingest, "_extract_pptx", return_value="pptx-text") as mocked:
            result = doc_ingest.extract_text("a.pptx", b"data")
        mocked.assert_called_once_with(b"data")
        self.assertEqual(result, "pptx-text")


class TestPlainTextExtraction(unittest.TestCase):
    """12日目追記: 📷ボタン統合に伴い追加したコード/テキスト系ファイル対応の確認。"""

    def test_py_file_dispatches_to_plain_text_extractor(self):
        with patch.object(doc_ingest, "_extract_plain_text", return_value="text") as mocked:
            result = doc_ingest.extract_text("a.py", b"data")
        mocked.assert_called_once_with(b"data")
        self.assertEqual(result, "text")

    def test_utf8_source_code_round_trips(self):
        code = 'print("こんにちは")\n'
        text = doc_ingest.extract_text("hello.py", code.encode("utf-8"))
        self.assertEqual(text, code)

    def test_json_csv_and_other_code_extensions_are_supported(self):
        for filename in ("data.json", "table.csv", "notes.md", "app.js", "query.sql"):
            with self.subTest(filename=filename):
                text = doc_ingest.extract_text(filename, "hello".encode("utf-8"))
                self.assertEqual(text, "hello")

    def test_cp932_source_falls_back_correctly(self):
        # メモ帳等でShift_JIS(CP932)保存されたテキストでも文字化けせずに読めること。
        text = "日本語のコメント".encode("cp932")
        result = doc_ingest.extract_text("memo.txt", text)
        self.assertEqual(result, "日本語のコメント")

    def test_image_extension_still_unsupported(self):
        # 画像はこのモジュールの対象外のまま(送信するとDEEPへ強制ルーティングされる
        # 別経路(pendingImages)であり、ナレッジへの永続登録はしない)。
        with self.assertRaises(doc_ingest.UnsupportedFileTypeError):
            doc_ingest.extract_text("photo.png", b"dummy")


class TestDocxExtraction(unittest.TestCase):
    """python-docxで実際にファイルを作成し、往復でテキストが取れることを確認する。"""

    def test_paragraphs_and_table_are_extracted(self):
        import docx

        document = docx.Document()
        document.add_paragraph("これは段落1です。")
        document.add_paragraph("これは段落2です。")
        table = document.add_table(rows=1, cols=2)
        table.rows[0].cells[0].text = "見出しA"
        table.rows[0].cells[1].text = "見出しB"

        buf = io.BytesIO()
        document.save(buf)

        text = doc_ingest.extract_text("test.docx", buf.getvalue())
        self.assertIn("これは段落1です。", text)
        self.assertIn("これは段落2です。", text)
        self.assertIn("見出しA", text)
        self.assertIn("見出しB", text)


class TestXlsxExtraction(unittest.TestCase):
    """openpyxlで実際にファイルを作成し、往復でテキストが取れることを確認する。"""

    def test_cells_are_extracted(self):
        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws.append(["名前", "値"])
        ws.append(["テスト太郎", 123])

        buf = io.BytesIO()
        wb.save(buf)

        text = doc_ingest.extract_text("test.xlsx", buf.getvalue())
        self.assertIn("Sheet1", text)
        self.assertIn("名前", text)
        self.assertIn("テスト太郎", text)
        self.assertIn("123", text)


class TestPptxExtraction(unittest.TestCase):
    """python-pptxで実際にファイルを作成し、往復でテキストが取れることを確認する。"""

    def test_slide_text_is_extracted(self):
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 白紙レイアウト
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        textbox.text_frame.text = "スライドの本文テキスト"

        buf = io.BytesIO()
        prs.save(buf)

        text = doc_ingest.extract_text("test.pptx", buf.getvalue())
        self.assertIn("slide 1", text)
        self.assertIn("スライドの本文テキスト", text)


class TestGetDocumentText(unittest.TestCase):
    """14日目①: get_document_text()。memory_storeをフェイクへ差し替え、LanceDBに
    触れずに「チャンクを連結して復元する」ロジック(方針1)だけを検証する。
    """

    def _install_fake_memory_store(self, rows):
        import pandas as pd

        class FakeTable:
            def __init__(self, rows):
                self._df = pd.DataFrame(rows)

            def count_rows(self):
                return len(self._df)

            def to_pandas(self):
                return self._df

        fake_module = type(sys)("memory_store")
        fake_module._table = lambda: FakeTable(rows)
        patcher = patch.dict(sys.modules, {"memory_store": fake_module})
        patcher.start()
        self.addCleanup(patcher.stop)

    def test_concatenates_chunks_for_registered_file(self):
        self._install_fake_memory_store(
            [
                {"source": "doc:a.pdf", "content": "1つ目のチャンク"},
                {"source": "doc:a.pdf", "content": "2つ目のチャンク"},
                {"source": "doc:other.pdf", "content": "無関係のチャンク"},
            ]
        )
        text = doc_ingest.get_document_text("a.pdf")
        self.assertEqual(text, "1つ目のチャンク\n\n2つ目のチャンク")

    def test_returns_none_for_unregistered_file(self):
        self._install_fake_memory_store([{"source": "doc:a.pdf", "content": "x"}])
        self.assertIsNone(doc_ingest.get_document_text("unknown.pdf"))

    def test_returns_none_when_table_is_empty(self):
        self._install_fake_memory_store([])
        self.assertIsNone(doc_ingest.get_document_text("a.pdf"))


if __name__ == "__main__":
    unittest.main()
